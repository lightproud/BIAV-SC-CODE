"""
Parser for Office and PDF files (.docx, .xlsx, .pptx, .pdf).

Each parser gracefully handles ImportError so the server starts
even when some optional libraries are not installed.

Chunking strategy:
  - .docx: chunk by paragraph
  - .xlsx: chunk by sheet, then by row groups (50 rows per chunk)
  - .pptx: chunk by slide
  - .pdf:  chunk by page
"""

from pathlib import Path
from typing import Any, Dict, List

from . import ParseResult

# Row group size for Excel chunking.
_XLSX_CHUNK_ROWS = 50


# -- DOCX parser ------------------------------------------------------------

def _parse_docx(file_path: str) -> ParseResult:
    """Parse a .docx file using python-docx."""
    try:
        from docx import Document
    except ImportError:
        return ParseResult(
            text="",
            metadata={
                "file": file_path,
                "error": "python-docx not installed. Run: pip install python-docx",
            },
            chunks=[],
        )

    doc = Document(file_path)
    paragraphs: List[str] = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            paragraphs.append(text)

    full_text = "\n\n".join(paragraphs)

    metadata: Dict[str, Any] = {
        "file": file_path,
        "extension": ".docx",
        "format": "docx",
        "paragraphs": len(paragraphs),
        "chunk_count": len(paragraphs),
    }

    return ParseResult(
        text=full_text,
        metadata=metadata,
        chunks=paragraphs,
    )


# -- XLSX parser -------------------------------------------------------------

def _parse_xlsx(file_path: str) -> ParseResult:
    """Parse a .xlsx file using openpyxl."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        return ParseResult(
            text="",
            metadata={
                "file": file_path,
                "error": "openpyxl not installed. Run: pip install openpyxl",
            },
            chunks=[],
        )

    wb = load_workbook(file_path, read_only=True, data_only=True)
    chunks: List[str] = []
    all_text_parts: List[str] = []
    total_rows = 0

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows: List[List[str]] = []
        for row in ws.iter_rows(values_only=True):
            cell_values = [str(cell) if cell is not None else "" for cell in row]
            rows.append(cell_values)
        total_rows += len(rows)

        # Chunk by row groups
        if not rows:
            continue

        header = rows[0] if rows else []
        for start in range(0, len(rows), _XLSX_CHUNK_ROWS):
            end = min(start + _XLSX_CHUNK_ROWS, len(rows))
            batch = rows[start:end]

            lines = []
            if start > 0 and header:
                # Repeat header for context in non-first chunks
                lines.append(f"[Sheet: {sheet_name}] " + " | ".join(header))
            elif header and start == 0:
                lines.append(f"[Sheet: {sheet_name}] " + " | ".join(batch[0]))
                batch = batch[1:]

            for row_cells in batch:
                lines.append(" | ".join(row_cells))

            chunk_text = "\n".join(lines)
            if chunk_text.strip():
                chunks.append(chunk_text)
                all_text_parts.append(chunk_text)

    wb.close()

    metadata: Dict[str, Any] = {
        "file": file_path,
        "extension": ".xlsx",
        "format": "xlsx",
        "sheets": wb.sheetnames if hasattr(wb, "sheetnames") else [],
        "total_rows": total_rows,
        "chunk_count": len(chunks),
    }

    return ParseResult(
        text="\n\n".join(all_text_parts),
        metadata=metadata,
        chunks=chunks,
    )


# -- PPTX parser ------------------------------------------------------------

def _parse_pptx(file_path: str) -> ParseResult:
    """Parse a .pptx file using python-pptx."""
    try:
        from pptx import Presentation
    except ImportError:
        return ParseResult(
            text="",
            metadata={
                "file": file_path,
                "error": "python-pptx not installed. Run: pip install python-pptx",
            },
            chunks=[],
        )

    prs = Presentation(file_path)
    chunks: List[str] = []

    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_texts: List[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_texts.append(text)

        if slide_texts:
            chunk = f"[Slide {slide_idx}]\n" + "\n".join(slide_texts)
            chunks.append(chunk)

    full_text = "\n\n".join(chunks)

    metadata: Dict[str, Any] = {
        "file": file_path,
        "extension": ".pptx",
        "format": "pptx",
        "slides": len(prs.slides),
        "chunk_count": len(chunks),
    }

    return ParseResult(
        text=full_text,
        metadata=metadata,
        chunks=chunks,
    )


# -- PDF parser --------------------------------------------------------------

def _parse_pdf(file_path: str) -> ParseResult:
    """Parse a .pdf file using PyMuPDF (fitz)."""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        return ParseResult(
            text="",
            metadata={
                "file": file_path,
                "error": "PyMuPDF not installed. Run: pip install pymupdf",
            },
            chunks=[],
        )

    doc = fitz.open(file_path)
    chunks: List[str] = []

    for page_idx in range(len(doc)):
        page = doc[page_idx]
        text = page.get_text().strip()
        if text:
            chunk = f"[Page {page_idx + 1}]\n{text}"
            chunks.append(chunk)

    doc.close()

    full_text = "\n\n".join(chunks)

    metadata: Dict[str, Any] = {
        "file": file_path,
        "extension": ".pdf",
        "format": "pdf",
        "pages": len(doc) if hasattr(doc, "__len__") else 0,
        "chunk_count": len(chunks),
    }

    return ParseResult(
        text=full_text,
        metadata=metadata,
        chunks=chunks,
    )


# -- Unified dispatcher ------------------------------------------------------

_OFFICE_PARSERS = {
    ".docx": _parse_docx,
    ".xlsx": _parse_xlsx,
    ".pptx": _parse_pptx,
    ".pdf": _parse_pdf,
}


def parse_office(file_path: str) -> ParseResult:
    """
    Parse an Office or PDF file.

    Args:
        file_path: Path to the file.

    Returns:
        ParseResult with extracted text, metadata, and chunks.
    """
    ext = Path(file_path).suffix.lower()
    parser_fn = _OFFICE_PARSERS.get(ext)

    if parser_fn is None:
        raise ValueError(f"No office parser for extension: {ext}")

    return parser_fn(file_path)
